from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT_DIR / "output"
PPTX_PATH = ROOT_DIR / "Packaging_Label_Defect_Inspection_Midterm.pptx"
PPTX_FALLBACK_PATH = ROOT_DIR / "Packaging_Label_Defect_Inspection_Midterm_updated.pptx"
BG = RGBColor(24, 28, 36)
PANEL = RGBColor(36, 42, 54)
ACCENT = RGBColor(255, 136, 0)
GREEN = RGBColor(72, 201, 176)
RED = RGBColor(231, 76, 60)
BLUE = RGBColor(86, 159, 255)
TEXT = RGBColor(245, 247, 250)
MUTED = RGBColor(191, 197, 207)


def read_results() -> list[dict[str, str]]:
    with (OUTPUT_DIR / "inspection_results.csv").open("r", newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.0), Inches(11.5), Inches(0.45))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED


def add_panel(slide, left: float, top: float, width: float, height: float):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL
    return shape


def add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str], size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.bullet = True


def add_caption(slide, left: float, top: float, width: float, text: str, color: RGBColor = MUTED, size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color


def add_metric_box(slide, left: float, top: float, width: float, height: float, value: str, label: str, accent: RGBColor) -> None:
    box = add_panel(slide, left, top, width, height)
    box.line.color.rgb = accent
    text = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(height - 0.2))
    tf = text.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = MUTED


def pick_row(rows: list[dict[str, str]], *, expected_status: str | None = None, defect_type: str | None = None, exclude_labels: set[str] | None = None) -> dict[str, str]:
    for row in rows:
        if expected_status is not None and row["expected_status"] != expected_status:
            continue
        if defect_type is not None and row["defect_type"] != defect_type:
            continue
        if exclude_labels and row["label_slug"] in exclude_labels:
            continue
        return row
    raise RuntimeError("Could not find a matching result row for the presentation.")


def summary_path(row: dict[str, str]) -> Path:
    return OUTPUT_DIR / "summaries" / f"{Path(row['filename']).stem}_summary.png"


def matches_path(row: dict[str, str]) -> Path:
    return OUTPUT_DIR / "matches" / f"{Path(row['filename']).stem}_matches.png"


def create_presentation() -> Path:
    results = read_results()
    total = len(results)
    correct = sum(row["correct"] == "True" for row in results)
    accuracy = correct / total if total else 0.0
    label_count = len({row["label_slug"] for row in results})
    pass_count = sum(row["expected_status"] == "PASS" for row in results)
    fail_count = sum(row["expected_status"] == "FAIL" for row in results)
    mean_inliers = sum(int(row["inliers"]) for row in results) / total if total else 0.0

    pass_row = pick_row(results, expected_status="PASS")
    fail_missing = pick_row(results, expected_status="FAIL", defect_type="missing_print", exclude_labels={pass_row["label_slug"]})
    fail_scratch = pick_row(results, expected_status="FAIL", defect_type="scratch", exclude_labels={pass_row["label_slug"], fail_missing["label_slug"]})
    gallery_img = OUTPUT_DIR / "reference_gallery.png"
    overview_img = OUTPUT_DIR / "demo_overview.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Multi-Label Packaging Defect Inspection", "Midterm project - Industrial Computer Vision Practice")
    add_bullets(
        slide,
        0.8,
        1.6,
        5.2,
        3.5,
        [
            "Goal: inspect different consumer-package fronts and decide PASS or FAIL automatically.",
            "Dataset: 5 real package labels collected from Open Food Facts product-front images.",
            "Course concepts used: SIFT matching, image alignment, thresholding, morphology, connected components.",
        ],
        18,
    )
    slide.shapes.add_picture(str(summary_path(fail_scratch)), Inches(6.25), Inches(1.45), width=Inches(6.35))
    add_caption(slide, 6.35, 6.9, 6.1, "Representative FAIL case with aligned comparison, defect mask, and final decision")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Realistic Example Data", "Five different product labels and four test images per label")
    add_panel(slide, 0.55, 1.35, 12.2, 5.75)
    slide.shapes.add_picture(str(gallery_img), Inches(0.8), Inches(1.7), width=Inches(11.7))
    add_bullets(slide, 0.95, 5.75, 11.0, 0.85, ["Labels used: cookies, cereal, chips, tea, and chocolate. Each label has 2 PASS samples and 2 FAIL samples with different defects."], 16)
    add_caption(slide, 0.9, 6.78, 11.4, "Source images: Open Food Facts package-front photos, then transformed into inspection samples")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Inspection Pipeline", "Reference-specific alignment first, defect localization second")
    steps = [("1", "Choose matching\nreference label"), ("2", "SIFT keypoints\nand matching"), ("3", "Homography\nalignment"), ("4", "Absdiff +\nthreshold"), ("5", "Morphology +\nconnected regions"), ("6", "PASS / FAIL")]
    left = 0.5
    for idx, (num, label) in enumerate(steps):
        box = add_panel(slide, left + idx * 2.08, 1.55, 1.78, 1.12)
        box.line.color.rgb = ACCENT if idx in (1, 2, 3) else BLUE
        text = slide.shapes.add_textbox(Inches(left + idx * 2.08 + 0.08), Inches(1.67), Inches(1.62), Inches(0.9))
        tf = text.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = num
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT
    slide.shapes.add_picture(str(matches_path(fail_missing)), Inches(0.95), Inches(3.15), width=Inches(5.25))
    add_bullets(slide, 6.65, 3.2, 5.3, 2.8, ["A separate reference is used for each label category.", "Alignment handles rotation, translation, and scale before comparison.", "Large connected residual regions are interpreted as label defects."], 18)
    add_caption(slide, 1.1, 6.55, 5.0, f"Example SIFT matches for {fail_missing['label_name']}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Representative Results", "One PASS case and two FAIL cases from different labels")
    slide.shapes.add_picture(str(summary_path(pass_row)), Inches(0.4), Inches(1.45), width=Inches(4.15))
    slide.shapes.add_picture(str(summary_path(fail_missing)), Inches(4.6), Inches(1.45), width=Inches(4.15))
    slide.shapes.add_picture(str(summary_path(fail_scratch)), Inches(8.8), Inches(1.45), width=Inches(4.15))
    add_caption(slide, 0.55, 6.75, 3.9, f"PASS: {pass_row['label_name']}", GREEN, 12)
    add_caption(slide, 4.75, 6.75, 3.9, f"FAIL: {fail_missing['label_name']} ({fail_missing['defect_type']})", RED, 12)
    add_caption(slide, 8.95, 6.75, 3.9, f"FAIL: {fail_scratch['label_name']} ({fail_scratch['defect_type']})", RED, 12)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Results Summary", "The same pipeline separates normal and defective samples across five labels")
    slide.shapes.add_picture(str(overview_img), Inches(0.55), Inches(1.5), width=Inches(7.4))
    add_metric_box(slide, 8.35, 1.65, 1.6, 1.0, str(label_count), "Labels", BLUE)
    add_metric_box(slide, 10.1, 1.65, 1.6, 1.0, str(total), "Samples", ACCENT)
    add_metric_box(slide, 8.35, 2.85, 1.6, 1.0, str(pass_count), "PASS samples", GREEN)
    add_metric_box(slide, 10.1, 2.85, 1.6, 1.0, str(fail_count), "FAIL samples", RED)
    add_metric_box(slide, 8.35, 4.05, 1.6, 1.0, f"{accuracy:.0%}", "Accuracy", GREEN)
    add_metric_box(slide, 10.1, 4.05, 1.6, 1.0, f"{mean_inliers:.0f}", "Avg. inliers", BLUE)
    add_bullets(slide, 8.25, 5.35, 3.9, 1.2, ["PASS and FAIL groups are clearly separated by defect area ratio.", "The threshold is fixed, not tuned separately for each label."], 15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_title(slide, "Conclusion", "Practical, explainable, and easy to present in 5 minutes")
    add_panel(slide, 0.7, 1.4, 5.9, 4.95)
    add_bullets(
        slide,
        0.95,
        1.72,
        5.25,
        4.2,
        [
            "Practical: packaging inspection is a real quality-control use case.",
            "Diverse: one algorithm handles 5 different labels instead of one synthetic box.",
            "Explainable: every decision is visible through matches, alignment, diff image, and mask.",
            "Current limitation: defects are simulated on real labels, not collected from a factory line.",
            "Next step: capture real phone or conveyor images and test lighting robustness.",
        ],
        17,
    )
    slide.shapes.add_picture(str(summary_path(fail_missing)), Inches(7.0), Inches(1.75), width=Inches(5.6))
    add_caption(slide, 7.1, 6.68, 5.3, "Online-source labels + explainable CV pipeline = strong midterm scope")

    try:
        prs.save(PPTX_PATH)
        output_path = PPTX_PATH
    except PermissionError:
        prs.save(PPTX_FALLBACK_PATH)
        output_path = PPTX_FALLBACK_PATH
    print(f"Saved presentation: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
